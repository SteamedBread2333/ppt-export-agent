from __future__ import annotations

import json
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from ppt_expert import AgentConfig, HostRuntime, create_ppt_agent
from ppt_expert.demo_runtime import fake_critique_images, fake_structured_generate
from ppt_expert.models import QualityIssue, StoryDesignBundle, VisionCritique


def _runtime(structured_generate=fake_structured_generate, critique_images=fake_critique_images):
    return HostRuntime(structured_generate=structured_generate, critique_images=critique_images)


def _config(tmp_path: Path) -> AgentConfig:
    return AgentConfig(
        output_root=tmp_path / "outputs",
        checkpoint_path=tmp_path / "checkpoints.sqlite3",
    )


async def _complete(agent, pending):
    result = pending
    while result["status"] == "interrupted":
        kind = result["request"]["type"]
        if kind == "intent_confirmation":
            result = await agent.resume(result["thread_id"], {"action": "continue"})
        elif kind == "recipe_confirmation":
            result = await agent.resume(result["thread_id"], {"action": "use"})
        elif kind == "delivery_confirmation":
            result = await agent.resume(result["thread_id"], {"action": "approve"})
        else:
            raise AssertionError(kind)
    return result


@pytest.mark.asyncio
async def test_end_to_end_interrupt_and_resume(tmp_path: Path) -> None:
    config = _config(tmp_path)
    runtime = _runtime()

    async with create_ppt_agent(runtime, config) as agent:
        interrupted = await agent.start("Create a PPT Expert workflow overview", project_name="demo")
        assert interrupted["status"] == "interrupted"
        assert interrupted["request"]["type"] == "recipe_confirmation"
        options = interrupted["request"]["options"]
        assert interrupted["request"]["recommended"] == options[0]["id"]
        assert options[0]["recommended"] is True
        assert {item["id"] for item in options} == {
            "consulting",
            "work_report",
            "civic",
            "art_market",
            "editorial",
            "history",
            "open",
        }
        assert interrupted["request"]["reason"]
        assert interrupted["request"]["recipe_id"]
        assert interrupted["request"]["palette"]

        pending = await agent.resume(interrupted["thread_id"], {"action": "use"})
        assert pending["status"] == "interrupted"
        assert pending["request"]["type"] == "delivery_confirmation"
        assert pending["request"]["validation"]["valid"] is True

        completed = await agent.resume(pending["thread_id"], {"action": "approve"})
        assert completed["status"] == "completed"
        assert completed["validation"]["valid"] is True

    artifacts = completed["artifacts"]
    assert Path(artifacts["story_path"]).exists()
    assert Path(artifacts["design_path"]).exists()
    assert Path(artifacts["report_path"]).exists()
    assert Path(artifacts["delivery_path"]).exists()
    assert (Path(artifacts["project_dir"]) / "tokens.json").exists()
    assert (Path(artifacts["project_dir"]) / "environment.json").exists()
    assert (Path(artifacts["project_dir"]) / "guards.json").exists()
    env = json.loads((Path(artifacts["project_dir"]) / "environment.json").read_text(encoding="utf-8"))
    assert env["visual_review"] == "full"
    assert Path(artifacts["montage_path"]).is_file()
    metrics = Path(artifacts["project_dir"]) / "metrics.jsonl"
    assert metrics.exists()
    presentation = Presentation(artifacts["pptx_path"])
    assert len(presentation.slides) == 4
    assert any(shape.has_chart for shape in presentation.slides[1].shapes)


@pytest.mark.asyncio
async def test_runtime_uses_host_model_structured_output() -> None:
    class Runnable:
        async def ainvoke(self, prompt):
            return fake_structured_generate(prompt, self.schema)

    class HostModel:
        def with_structured_output(self, schema):
            runnable = Runnable()
            runnable.schema = schema
            return runnable

    from ppt_expert.models import OutlinePlan

    runtime = HostRuntime(model=HostModel())
    result = await runtime.generate_structured("test", OutlinePlan)
    assert result.title == "PPT Expert Demo"


@pytest.mark.asyncio
async def test_validation_failure_routes_through_repair(tmp_path: Path) -> None:
    story_calls = 0

    def host_generate(prompt, schema):
        nonlocal story_calls
        result = fake_structured_generate(prompt, schema)
        if schema is StoryDesignBundle:
            story_calls += 1
            if story_calls == 1:
                result = result.model_copy(deep=True)
                result.pages[0].content = ["x" * 400]
        return result

    config = _config(tmp_path)
    async with create_ppt_agent(_runtime(structured_generate=host_generate), config) as agent:
        interrupted = await agent.start("repair test")
        completed = await _complete(agent, interrupted)
        state = await agent.state(interrupted["thread_id"])

    assert completed["validation"]["valid"] is True
    assert state["repair_attempts"] >= 1
    assert story_calls == 2


@pytest.mark.asyncio
async def test_template_still_renders(tmp_path: Path) -> None:
    template = tmp_path / "template.pptx"
    source = Presentation()
    source.slide_width = Inches(10)
    source.slide_height = Inches(7.5)
    source.slides.add_slide(source.slide_layouts[6])
    source.save(template)
    config = _config(tmp_path)
    async with create_ppt_agent(_runtime(), config) as agent:
        pending = await agent.start("template test", template_path=template)
        completed = await _complete(agent, pending)

    presentation = Presentation(completed["artifacts"]["pptx_path"])
    assert len(presentation.slides) == 4
    assert presentation.slide_width == Inches(10)
    assert completed["validation"]["valid"] is True
    assert Path(completed["artifacts"]["montage_path"]).is_file()
    marker = Path(completed["artifacts"]["project_dir"]) / "template.json"
    assert marker.exists()


@pytest.mark.asyncio
async def test_reference_images_do_not_block_recipe_gate(tmp_path: Path) -> None:
    reference = tmp_path / "reference.png"
    Image.new("RGB", (100, 100), "#336699").save(reference)
    config = _config(tmp_path)
    async with create_ppt_agent(_runtime(), config) as agent:
        pending = await agent.start("reference test", reference_images=[reference])
        assert pending["request"]["type"] == "recipe_confirmation"
        assert pending["request"]["options"]
        completed = await _complete(agent, pending)
    assert completed["status"] == "completed"


@pytest.mark.asyncio
async def test_user_can_override_the_recommended_recipe(tmp_path: Path) -> None:
    config = _config(tmp_path)
    async with create_ppt_agent(_runtime(), config) as agent:
        pending = await agent.start("Create a PPT Expert workflow overview", project_name="override")
        recommended = pending["request"]["recommended"]
        assert recommended
        pending = await agent.resume(pending["thread_id"], {"action": "history"})
        assert pending["request"]["type"] == "delivery_confirmation"
        state = await agent.state(pending["thread_id"])
        assert state["recipe_id"] == "history"
        assert recommended != "history"
        completed = await agent.resume(pending["thread_id"], {"action": "approve"})
    assert completed["status"] == "completed"


@pytest.mark.asyncio
async def test_montage_critique_repairs_flagged_pages(tmp_path: Path) -> None:
    critiques = {"n": 0}

    def critic(prompt, paths, schema):
        del prompt, schema
        assert Path(paths[0]).is_file()
        critiques["n"] += 1
        if critiques["n"] == 1:
            return VisionCritique(
                score=40,
                issues=[
                    QualityIssue(
                        code="montage_density",
                        message="Page 3 is too dense on the contact sheet",
                        page=3,
                        severity="error",
                        cause="The montage shows colliding modules",
                        repair_scope="slide",
                        acceptance="Density is resolved on rebuild",
                    )
                ],
                notes="first pass",
            )
        return VisionCritique(score=88, issues=[], notes="cleared")

    story_calls = 0

    def host_generate(prompt, schema):
        nonlocal story_calls
        result = fake_structured_generate(prompt, schema)
        if schema is StoryDesignBundle:
            story_calls += 1
        return result

    async with create_ppt_agent(
        _runtime(structured_generate=host_generate, critique_images=critic),
        _config(tmp_path),
    ) as agent:
        interrupted = await agent.start("montage repair")
        completed = await _complete(agent, interrupted)
        state = await agent.state(interrupted["thread_id"])

    assert Path(completed["artifacts"]["montage_path"]).is_file()
    assert critiques["n"] >= 2
    assert story_calls >= 2
    assert state["repair_attempts"] >= 1
