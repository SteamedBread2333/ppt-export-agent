from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from ppt_expert.config import AgentConfig
from ppt_expert.models import (
    ChartSeries,
    ChartSpec,
    ChartType,
    DesignSpec,
    LayoutType,
    OutlinePage,
    OutlinePlan,
    PageRole,
    RecipeId,
    SlideFamily,
    StoryPage,
)
from ppt_expert.pptx import render_presentation
from ppt_expert.recipes import tokens_for
from ppt_expert.review import review_volume
from ppt_expert.validation import validate_presentation


def _design() -> DesignSpec:
    return DesignSpec(
        style_name="consulting",
        mood="Clear and composed",
        primary="#1F4E79",
        secondary="#44505C",
        background="#F7F8FA",
        text="#1B242C",
        accent="#1F4E79",
        illustration_style="vector_first",
        typography_profile="consulting",
        token_palette=list(tokens_for(RecipeId.CONSULTING).palette_hex()),
    )


def _page(
    number: int,
    title: str,
    content: list[str],
    *,
    role: PageRole,
    chart: ChartSpec | None = None,
) -> StoryPage:
    return StoryPage(
        number=number,
        title=title,
        content=content,
        visual_direction=title,
        layout=LayoutType.TEXT,
        family=SlideFamily.COVER if role == PageRole.COVER else SlideFamily.CHART_INTERPRETATION,
        role=role,
        chart=chart,
        speaker_notes=title,
    )


def test_native_edit_keeps_template_art_and_size(tmp_path: Path) -> None:
    template = tmp_path / "brand.pptx"
    source = Presentation()
    source.slide_width = Inches(10)
    source.slide_height = Inches(7.5)
    slide = source.slides.add_slide(source.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(7), Inches(1))
    title.text_frame.text = "Old cover title"
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(0.3), Inches(1.5), Inches(0.8)
    )
    chip.name = "brand-chip"
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(0xC4, 0x5C, 0x26)
    source.save(template)

    pages = [
        _page(1, "Keep the brand chip", ["The template stays the canvas"], role=PageRole.COVER),
        _page(2, "Clone when the outline is longer", ["One template page becomes two"], role=PageRole.CLOSE),
    ]
    outline = OutlinePlan(
        title="Native edit",
        pages=[OutlinePage(number=page.number, title=page.title, core_content=page.content) for page in pages],
    )
    path = tmp_path / "edited.pptx"
    render_presentation(pages, _design(), {}, path, AgentConfig(), template_path=template)
    presentation = Presentation(path)

    assert presentation.slide_width == Inches(10)
    assert presentation.slide_height == Inches(7.5)
    assert len(presentation.slides) == 2
    names = [shape.name for slide in presentation.slides for shape in slide.shapes]
    assert names.count("brand-chip") == 2
    texts = [
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]
    assert any("Keep the brand chip" in text for text in texts)
    assert any("The template stays the canvas" in text for text in texts)

    report = validate_presentation(path, outline, pages, _design(), {}, native_edit=True)
    assert report.valid, [issue.model_dump() for issue in report.issues]
    assert not any(issue.code in {"palette_violation", "font_violation"} for issue in report.issues)

    review = review_volume(path, pages, tmp_path, native_edit=True)
    assert not any(issue.code in {"card_soup", "cramped_header", "empty_bottom"} for issue in review.issues)


def test_native_edit_assigns_chart_pages_to_chart_slides(tmp_path: Path) -> None:
    template = tmp_path / "charts.pptx"
    source = Presentation()
    source.slide_width = Inches(10)
    source.slide_height = Inches(7.5)
    cover = source.slides.add_slide(source.slide_layouts[6])
    heading = cover.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(1))
    heading.text_frame.text = "Cover slot"
    chart_slide = source.slides.add_slide(source.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("Old", (1.0, 2.0))
    chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(7), Inches(4), data
    )
    source.save(template)

    pages = [
        _page(1, "Cover stays on the first slot", ["No chart on cover"], role=PageRole.COVER),
        _page(
            2,
            "Revisions inflect before price",
            ["Follow the slope, not the level"],
            role=PageRole.CONTEXT,
            chart=ChartSpec(
                chart_type=ChartType.COLUMN,
                categories=["Q1", "Q2"],
                series=[ChartSeries(name="Focus", values=[8, 12])],
            ),
        ),
    ]
    path = tmp_path / "chart-edit.pptx"
    render_presentation(pages, _design(), {}, path, AgentConfig(), template_path=template)
    presentation = Presentation(path)

    assert len(presentation.slides) == 2
    cover_has_chart = any(
        getattr(shape, "has_chart", False) for shape in presentation.slides[0].shapes
    )
    body_charts = [
        shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False)
    ]
    assert cover_has_chart is False
    assert body_charts
    categories = [str(category) for category in body_charts[0].plots[0].categories]
    assert categories == ["Q1", "Q2"]
