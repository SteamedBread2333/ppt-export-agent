from pathlib import Path

from pptx import Presentation

from ppt_expert.config import AgentConfig
from ppt_expert.guards import inspect_guards
from ppt_expert.models import KPIItem, LayoutType, PageRole, RecipeId, SlideFamily, StoryPage
from ppt_expert.pptx import render_presentation
from ppt_expert.recipes import tokens_for


def test_short_numeric_tokens_disable_wrap(tmp_path: Path) -> None:
    page = StoryPage(
        number=1,
        title="Inflection, not a linear rebound",
        content=["Own the turn"],
        visual_direction="Cover",
        layout=LayoutType.HERO,
        family=SlideFamily.COVER,
        role=PageRole.COVER,
        kpis=[KPIItem(value="36.8%", label="Consensus upside", note="CSI 300")],
        speaker_notes="Cover intent",
    )
    path = tmp_path / "token.pptx"
    tokens = tokens_for(RecipeId.CONSULTING)
    render_presentation([page], tokens.to_design_spec(), {}, path, AgentConfig(), tokens=tokens)
    wrapped = []
    for shape in Presentation(path).slides[0].shapes:
        if getattr(shape, "has_text_frame", False) and "36.8%" in shape.text:
            wrapped.append(shape.text_frame.word_wrap)
    assert wrapped
    assert not any(wrapped)
    assert inspect_guards(path).clean
