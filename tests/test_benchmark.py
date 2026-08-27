from pathlib import Path

from pptx import Presentation

from ppt_expert.benchmarks import STRATEGY_DESIGN, extra_analytical_pages, strategy_benchmark_pages
from ppt_expert.config import AgentConfig
from ppt_expert.models import OutlinePage, OutlinePlan
from ppt_expert.pptx import render_presentation
from ppt_expert.quality import score_deck
from ppt_expert.validation import validate_presentation


def test_strategy_benchmark_meets_acceptance(tmp_path: Path) -> None:
    pages = strategy_benchmark_pages()
    path = tmp_path / "strategy.pptx"
    render_presentation(pages, STRATEGY_DESIGN, {}, path, AgentConfig())
    outline = OutlinePlan(
        title="Strategy",
        pages=[
            OutlinePage(number=page.number, title=page.title, core_content=page.content)
            for page in pages
        ],
    )
    report = validate_presentation(path, outline, pages, STRATEGY_DESIGN, {})
    quality = score_deck(pages, STRATEGY_DESIGN)
    presentation = Presentation(path)
    charts = sum(1 for slide in presentation.slides for shape in slide.shapes if shape.has_chart)
    tables = sum(1 for slide in presentation.slides for shape in slide.shapes if shape.has_table)
    families = [page.resolved_family() for page in pages]

    assert report.valid is True, report.model_dump()
    assert quality.score >= 90
    assert all(value >= 82 for value in quality.dimensions.values())
    assert not quality.blocking_issues
    assert charts >= 3
    assert tables >= 1
    assert any(page.scenarios for page in pages)
    assert any(page.allocation for page in pages)
    assert all(page.source_note for page in pages if page.number in {1, 3, 5, 6, 7})
    assert not any(page.image_id for page in pages)
    assert len(set(families)) >= 6
    assert families[0].value == "cover"
    assert families[-1].value == "conclusion"


def test_waterfall_and_heatmap_render(tmp_path: Path) -> None:
    pages = extra_analytical_pages()
    path = tmp_path / "extra.pptx"
    render_presentation(pages, STRATEGY_DESIGN, {}, path, AgentConfig())
    outline = OutlinePlan(
        title="Extra",
        pages=[
            OutlinePage(number=page.number, title=page.title, core_content=page.content)
            for page in pages
        ],
    )
    report = validate_presentation(path, outline, pages, STRATEGY_DESIGN, {})
    assert report.valid is True, report.model_dump()
    assert any(shape.has_table for shape in Presentation(path).slides[1].shapes)
