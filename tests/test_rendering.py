import re
import zipfile
from io import BytesIO
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_expert.assets import generate_assets
from ppt_expert.config import AgentConfig
from ppt_expert.models import (
    ChartSeries,
    ChartSpec,
    ChartType,
    DesignSpec,
    ImageRequest,
    KPIItem,
    LayoutType,
    OutlinePage,
    OutlinePlan,
    PageRole,
    RecipeId,
    SlideFamily,
    StoryPage,
    TableSpec,
)
from ppt_expert.pptx import render_presentation
from ppt_expert.recipes import tokens_for
from ppt_expert.review import review_volume
from ppt_expert.runtime import HostRuntime
from ppt_expert.validation import validate_presentation


def _design() -> DesignSpec:
    return DesignSpec(
        style_name="consulting",
        mood="Clear and composed",
        primary="#1F4E79",
        secondary="#44505C",
        background="#F7F8FA",
        text="#1B242C",
        accent="#1F4E79",
        illustration_style="vector_first",
        typography_profile="consulting",
        token_palette=list(tokens_for(RecipeId.CONSULTING).palette_hex()),
    )


def test_every_layout_renders_and_validates(tmp_path: Path) -> None:
    layouts = list(LayoutType)
    pages = [
        StoryPage(
            number=index,
            title=f"Slide {index}",
            content=[f"Core message {index}"],
            visual_direction="Abstract visual composition",
            layout=layout,
        )
        for index, layout in enumerate(layouts, 1)
    ]
    outline = OutlinePlan(
        title="Layout Test",
        pages=[
            OutlinePage(number=page.number, title=page.title, core_content=page.content)
            for page in pages
        ],
    )
    path = tmp_path / "layouts.pptx"
    render_presentation(pages, _design(), {}, path, AgentConfig())
    report = validate_presentation(path, outline, pages, _design(), {})

    assert report.valid is True
    assert len(Presentation(path).slides) == len(layouts)


def test_native_chart_table_and_kpis_validate(tmp_path: Path) -> None:
    pages = [
        StoryPage(
            number=1,
            title="Cover with evidence",
            content=["Strategy remains overweight cyclical earnings"],
            visual_direction="KPI cover",
            layout=LayoutType.HERO,
            family=SlideFamily.COVER,
            eyebrow="2026H2",
            kpis=[
                KPIItem(value="36.8%", label="Consensus upside", note="CSI 300"),
                KPIItem(value="3", label="Scenarios", note="Base / bull / bear"),
            ],
        ),
        StoryPage(
            number=2,
            title="Earnings revision path",
            content=["Revisions inflect before price", "Follow the slope, not the level"],
            visual_direction="Line chart of revisions",
            layout=LayoutType.TEXT,
            family=SlideFamily.CHART_INTERPRETATION,
            takeaway="The turn is already in the data.",
            chart=ChartSpec(
                chart_type=ChartType.LINE,
                categories=["Q1", "Q2", "Q3", "Q4"],
                series=[ChartSeries(name="Revisions", values=[-1.2, -0.4, 0.6, 1.1])],
            ),
        ),
        StoryPage(
            number=3,
            title="Allocation",
            content=["Keep a barbell until breadth expands"],
            visual_direction="Allocation strip",
            layout=LayoutType.TEXT,
            family=SlideFamily.TABLE_COMPARISON,
            table=TableSpec(
                headers=["Sleeve", "Weight", "Bias"],
                rows=[["Cyclicals", "40%", "Overweight"], ["Quality", "35%", "Neutral"]],
                highlight_row=0,
            ),
        ),
    ]
    outline = OutlinePlan(
        title="Analytics",
        pages=[
            OutlinePage(number=page.number, title=page.title, core_content=page.content)
            for page in pages
        ],
    )
    path = tmp_path / "analytics.pptx"
    design = _design()
    render_presentation(pages, design, {}, path, AgentConfig())
    report = validate_presentation(path, outline, pages, design, {})
    presentation = Presentation(path)

    assert report.valid is True, report.model_dump()
    assert any(shape.has_chart for shape in presentation.slides[1].shapes)
    assert any(shape.has_table for shape in presentation.slides[2].shapes)
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        assert len([n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")]) == 3
        assert len([n for n in names if n.startswith("ppt/notesSlides/") and n.endswith(".xml")]) >= 3
        xml = " ".join(archive.read(n).decode("utf-8", errors="ignore") for n in names if n.endswith(".xml"))
        assert not re.search(r"undefined|NaN|\[object", xml)


def test_pillar_cards_use_hierarchy_instead_of_tall_empty_columns(tmp_path: Path) -> None:
    pages = [
        StoryPage(
            number=1,
            title="三条主线：出海制造、科技自主、红利重估",
            content=[
                "出海制造：海外收入 34%, 盈利 +14%; 风险：关税与汇率",
                "科技自主：设备订单 +42%, 盈利 +21%; 风险：估值抬升过快",
                "红利重估：股息率 4.8%, 增配空间 +1.6pct; 风险：利率反弹",
                "机制：三主线相关系数 <0.4, 等权配置即构成再平衡",
            ],
            visual_direction="Four theme modules",
            layout=LayoutType.DATA_CARDS,
            family=SlideFamily.PILLARS,
            role=PageRole.EXPANSION,
            eyebrow="行业配置",
            speaker_notes="Four expansion columns.",
        )
    ]
    outline = OutlinePlan(
        title="Pillars",
        pages=[OutlinePage(number=1, title=pages[0].title, core_content=pages[0].content)],
    )
    path = tmp_path / "pillars.pptx"
    design = _design()
    render_presentation(pages, design, {}, path, AgentConfig())
    report = validate_presentation(path, outline, pages, design, {})
    presentation = Presentation(path)
    cards = [
        shape
        for shape in presentation.slides[0].shapes
        if shape.has_text_frame and "海外收入" in shape.text
    ]

    assert report.valid is True, report.model_dump()
    assert any("01" in shape.text for shape in presentation.slides[0].shapes if shape.has_text_frame)
    assert cards
    assert cards[0].height < Inches(3.2)


def test_chart_interpretation_sizes_copy_to_content(tmp_path: Path) -> None:
    pages = [
        StoryPage(
            number=1,
            title="增长换挡而非失速，政策托底意愿强于弹性",
            content=[
                "内需：社零 +4.6%，弱修复延续",
                "外需：出口 +3.8%，韧性尚存",
                "政策：赤字率 4.0%，托底意愿强于弹性",
            ],
            visual_direction="PMI chart",
            layout=LayoutType.TEXT,
            family=SlideFamily.CHART_INTERPRETATION,
            role=PageRole.CONTEXT,
            eyebrow="宏观动能",
            takeaway="研判：分子端修复依靠盈利回升，而非宏观杠杆重新扩张",
            chart=ChartSpec(
                chart_type=ChartType.LINE,
                title="制造业 / 非制造业 PMI",
                categories=["24Q1", "24Q2", "25Q4"],
                series=[ChartSeries(name="制造业 PMI", values=[50.0, 49.2, 50.1])],
            ),
        )
    ]
    outline = OutlinePlan(
        title="Macro",
        pages=[OutlinePage(number=1, title=pages[0].title, core_content=pages[0].content)],
    )
    path = tmp_path / "chart.pptx"
    design = _design()
    render_presentation(pages, design, {}, path, AgentConfig())
    report = validate_presentation(path, outline, pages, design, {})
    takeaway = next(
        shape
        for shape in Presentation(path).slides[0].shapes
        if shape.has_text_frame and "分子端修复" in shape.text
    )

    assert report.valid is True, report.model_dump()
    assert takeaway.height < Inches(2.2)
    chart = next(shape.chart for shape in Presentation(path).slides[0].shapes if shape.has_chart)
    plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    sp_pr = plot_area.find(qn("c:spPr"))
    assert sp_pr is not None
    assert sp_pr.find(qn("a:noFill")) is not None
    cat_ax = plot_area.find(qn("c:catAx"))
    body_pr = cat_ax.find(qn("c:txPr")).find(qn("a:bodyPr"))
    assert body_pr.get("rot") == "0"


def test_header_rule_and_footer_are_locked(tmp_path: Path) -> None:
    pages = [
        StoryPage(
            number=1,
            title="Cover assertion for the briefing",
            content=["Own the inflection"],
            visual_direction="Cover",
            layout=LayoutType.HERO,
            family=SlideFamily.COVER,
            role=PageRole.COVER,
            kpis=[KPIItem(value="3", label="Scenarios")],
            speaker_notes="Cover",
        ),
        StoryPage(
            number=2,
            title="Three conditions still have to hold",
            content=["Breadth expands", "Rates stay orderly", "Credit turns"],
            visual_direction="Overview",
            layout=LayoutType.DATA_CARDS,
            family=SlideFamily.EXECUTIVE_SUMMARY,
            role=PageRole.OVERVIEW,
            speaker_notes="Overview",
        ),
        StoryPage(
            number=3,
            title="Process is the hedge against narrative",
            content=["Rebalance monthly", "Cut on revision failure", "Add only on breadth"],
            visual_direction="Expansion",
            layout=LayoutType.DATA_CARDS,
            family=SlideFamily.PILLARS,
            role=PageRole.EXPANSION,
            speaker_notes="Expansion",
        ),
        StoryPage(
            number=4,
            title="Close with a decision, not a slogan",
            content=["Act this quarter"],
            visual_direction="Close",
            layout=LayoutType.TEXT,
            family=SlideFamily.CONCLUSION,
            role=PageRole.CLOSE,
            speaker_notes="Close",
        ),
    ]
    path = tmp_path / "chrome.pptx"
    render_presentation(pages, _design(), {}, path, AgentConfig())
    presentation = Presentation(path)
    footer_tops = []
    header_tops = []
    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text in {"01", "02", "03", "04"} and round(shape.top.inches, 1) >= 7.0:
                footer_tops.append(round(shape.top.inches, 2))
            if slide_index in {1, 2} and pages[slide_index].title in text:
                header_tops.append(round(shape.top.inches, 2))
    assert footer_tops
    assert len(set(footer_tops)) == 1
    assert header_tops
    assert len(set(header_tops)) == 1
    overview_index_labels = [
        round(shape.top.inches, 2)
        for shape in presentation.slides[1].shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip() == "01"
        and shape.top.inches < 3
    ]
    assert overview_index_labels
    assert min(overview_index_labels) >= 1.45



@pytest.mark.asyncio
async def test_missing_host_image_tool_uses_palette_placeholder(tmp_path: Path) -> None:
    request = ImageRequest(image_id="hero", page_numbers=[1], prompt="test")
    paths = await generate_assets(HostRuntime(), [request], _design(), tmp_path)

    generated = Path(paths["hero"])
    assert generated.exists()
    with Image.open(generated) as image:
        assert image.size == (1536, 1024)


@pytest.mark.asyncio
async def test_jpeg_written_to_png_path_is_normalized(tmp_path: Path) -> None:
    def write_disguised_jpeg(request, output_path):
        image = Image.new("RGB", (64, 48), (220, 20, 30))
        image.save(output_path, format="JPEG")
        return output_path

    request = ImageRequest(image_id="hero", page_numbers=[1], prompt="jpeg bytes")
    runtime = HostRuntime(image_generate=write_disguised_jpeg)
    paths = await generate_assets(runtime, [request], _design(), tmp_path)

    generated = Path(paths["hero"])
    assert generated.suffix == ".png"
    with Image.open(generated) as image:
        assert image.format == "PNG"
        assert image.getpixel((32, 24))[0] > 180


@pytest.mark.asyncio
async def test_jpg_sibling_written_by_host_is_discovered_and_converted(tmp_path: Path) -> None:
    def write_jpg_sibling(request, output_path):
        jpeg_path = output_path.with_suffix(".jpg")
        buffer = BytesIO()
        Image.new("RGB", (80, 60), (15, 180, 60)).save(buffer, format="JPEG")
        jpeg_path.write_bytes(buffer.getvalue())

    request = ImageRequest(image_id="hero", page_numbers=[1], prompt="jpg sibling")
    runtime = HostRuntime(image_generate=write_jpg_sibling)
    paths = await generate_assets(runtime, [request], _design(), tmp_path)

    generated = Path(paths["hero"])
    assert generated.suffix == ".png"
    with Image.open(generated) as image:
        assert image.format == "PNG"
        assert image.getpixel((40, 30))[1] > 140


def test_compose_modules_use_role_tokens_not_stray_hex() -> None:
    allowed = {"#FFFFFF", "#000000"}
    root = Path(__file__).resolve().parents[1] / "src" / "ppt_expert" / "pptx"
    for name in ("slides.py", "primitives.py", "renderer.py"):
        found = set(re.findall(r"#[0-9A-Fa-f]{6}", (root / name).read_text(encoding="utf-8")))
        assert found <= allowed, f"{name} has stray hex: {found - allowed}"


def test_package_xml_has_notes_and_no_placeholders(tmp_path: Path) -> None:
    pages = [
        StoryPage(
            number=1,
            title="Cover assertion for the briefing",
            content=["Own the inflection"],
            visual_direction="Cover",
            layout=LayoutType.HERO,
            family=SlideFamily.COVER,
            role=PageRole.COVER,
            speaker_notes="Cover intent",
            kpis=[KPIItem(value="3", label="Scenarios")],
        )
    ]
    path = tmp_path / "xml.pptx"
    render_presentation(pages, _design(), {}, path, AgentConfig())
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        notes = [name for name in names if name.startswith("ppt/notesSlides/") and name.endswith(".xml")]
        blob = " ".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in names
            if name.endswith(".xml")
        )
    assert notes
    assert not re.search(r"undefined|NaN|\[object", blob)


def test_overview_and_context_fail_card_soup_review_if_tiled(tmp_path: Path) -> None:
    pages = [
        StoryPage(
            number=1,
            title="增长有底、盈利渐明：以杠铃结构穿越风格再平衡",
            content=[
                "增长换挡而非失速：名义增长 4.8%–5.2%",
                "盈利底部渐明：非金融由 -2.1% 回升至 +6.5%",
                "风格再平衡：拥挤度 96%→71%",
            ],
            visual_direction="Overview",
            layout=LayoutType.DATA_CARDS,
            family=SlideFamily.EXECUTIVE_SUMMARY,
            role=PageRole.OVERVIEW,
            takeaway="以杠铃结构穿越再平衡。",
            speaker_notes="Overview",
            kpis=[
                KPIItem(value="4,150–4,600", label="沪深300", note="目标区间"),
                KPIItem(value="+6.5%", label="盈利增速", note="2026E"),
            ],
        ),
        StoryPage(
            number=2,
            title="增长换挡而非失速，政策托底意愿强于弹性",
            content=["内需：社零 +4.6%", "外需：出口 +3.8%", "政策：赤字率 4.0%"],
            visual_direction="PMI",
            layout=LayoutType.TEXT,
            family=SlideFamily.CHART_INTERPRETATION,
            role=PageRole.CONTEXT,
            takeaway="分子端修复依靠盈利回升，而非宏观杠杆重新扩张",
            speaker_notes="Context",
            chart=ChartSpec(
                chart_type=ChartType.LINE,
                title="制造业 / 非制造业 PMI",
                categories=["24Q1", "24Q2", "25Q4"],
                series=[ChartSeries(name="制造业 PMI", values=[50.0, 49.2, 50.1])],
            ),
        ),
    ]
    path = tmp_path / "taste.pptx"
    render_presentation(pages, _design(), {}, path, AgentConfig())
    review = review_volume(path, pages, tmp_path, visual_review="degraded")
    codes = {issue.code for issue in review.issues}
    assert "card_soup" not in codes
    assert "cramped_header" not in codes
