from pathlib import Path

from pptx import Presentation

from ppt_expert.preview import render_montage, render_representative_pages


def test_render_montage_writes_contact_sheet(tmp_path: Path) -> None:
    pptx = tmp_path / "one.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(pptx)
    render_dir = tmp_path / "render"
    montage, pdf = render_montage(pptx, render_dir, dpi=70)
    assert Path(montage).is_file()
    assert Path(pdf).is_file()
    assert not list(render_dir.glob("pg*.png"))
    pages = render_representative_pages(pdf, render_dir, [1], dpi=72)
    assert pages
    assert all(Path(path).is_file() for path in pages)
