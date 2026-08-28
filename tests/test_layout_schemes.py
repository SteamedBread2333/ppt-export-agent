from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE

from ppt_expert.config import AgentConfig
from ppt_expert.models import (
    ChartSeries,
    ChartSpec,
    ChartType,
    KPIItem,
    LayoutType,
    PageRole,
    RecipeId,
    SlideFamily,
    StoryPage,
)
from ppt_expert.pptx import render_presentation
from ppt_expert.recipes import tokens_for
from ppt_expert.review import review_volume


def _overview(number: int, content: list[str], *, kpis: bool = False) -> StoryPage:
    return StoryPage(
        number=number,
        title=f"判断 {number} 仍然成立",
        content=content,
        visual_direction="Overview",
        layout=LayoutType.DATA_CARDS,
        family=SlideFamily.EXECUTIVE_SUMMARY,
        role=PageRole.OVERVIEW,
        takeaway="以杠铃结构穿越再平衡。",
        speaker_notes="Overview",
        kpis=[KPIItem(value="4,150", label="沪深300"), KPIItem(value="+6.5%", label="盈利")] if kpis else [],
    )


def _expansion(number: int) -> StoryPage:
    return StoryPage(
        number=number,
        title="四条主线",
        content=[
            "出海制造：海外收入 34%",
            "科技自主：设备订单 +42%",
            "红利重估：股息率 4.8%",
            "机制：相关系数 <0.4",
        ],
        visual_direction="Expansion",
        layout=LayoutType.DATA_CARDS,
        family=SlideFamily.PILLARS,
        role=PageRole.EXPANSION,
        speaker_notes="Expansion",
    )


def _context(number: int) -> StoryPage:
    return StoryPage(
        number=number,
        title="增长换挡而非失速",
        content=["内需：社零 +4.6%", "外需：出口 +3.8%", "政策：赤字率 4.0%"],
        visual_direction="PMI",
        layout=LayoutType.TEXT,
        family=SlideFamily.CHART_INTERPRETATION,
        role=PageRole.CONTEXT,
        takeaway="分子端修复依靠盈利回升",
        speaker_notes="Context",
        chart=ChartSpec(
            chart_type=ChartType.LINE,
            title="PMI",
            categories=["24Q1", "24Q2", "25Q4"],
            series=[ChartSeries(name="制造业 PMI", values=[50.0, 49.2, 50.1])],
        ),
    )


def _body_cross(slide) -> bool:
    vlines = []
    hlines = []
    for shape in slide.shapes:
        width = shape.width.inches
        height = shape.height.inches
        top = shape.top.inches
        if top < 1.32 or top > 6.2:
            continue
        if 0.006 <= width <= 0.03 and height > 1.0:
            vlines.append(shape)
        if 0.006 <= height <= 0.03 and width > 3.0:
            hlines.append(shape)
    return bool(vlines) and bool(hlines)


def _silhouette(slide) -> tuple:
    vlines = hlines = 0
    fill_tops: list[float] = []
    for shape in slide.shapes:
        width = shape.width.inches
        height = shape.height.inches
        top = shape.top.inches
        if top < 1.25 or top > 6.3:
            continue
        if 0.006 <= width <= 0.03 and height > 0.8:
            vlines += 1
        if 0.006 <= height <= 0.03 and width > 2.5:
            hlines += 1
        if width >= 1.4 and height >= 0.45:
            fill = getattr(shape, "fill", None)
            try:
                if fill is not None and fill.type == MSO_FILL_TYPE.SOLID:
                    fill_tops.append(round(top, 1))
            except (AttributeError, TypeError, ValueError):
                pass
    return (vlines, hlines, tuple(sorted(fill_tops)), len(slide.shapes))


def _render(recipe: RecipeId, pages: list[StoryPage], tmp_path: Path, name: str):
    tokens = tokens_for(recipe)
    path = tmp_path / f"{name}.pptx"
    render_presentation(pages, tokens.to_design_spec(), {}, path, AgentConfig(), tokens=tokens)
    return Presentation(path), path, tokens


def test_non_consulting_recipes_do_not_share_the_hairline_cross(tmp_path: Path) -> None:
    pages = [_overview(1, ["增长换挡", "盈利回升", "风格再平衡"], kpis=True)]
    for recipe in (
        RecipeId.WORK_REPORT,
        RecipeId.CIVIC,
        RecipeId.ART_MARKET,
        RecipeId.EDITORIAL,
        RecipeId.HISTORY,
        RecipeId.OPEN,
    ):
        slide, _, _ = _render(recipe, pages, tmp_path, recipe.value)
        assert not _body_cross(slide.slides[0]), recipe.value


def test_filled_schemes_are_allowed_to_use_surfaces(tmp_path: Path) -> None:
    pages = [_overview(1, ["增长换挡", "盈利回升", "风格再平衡"], kpis=True)]
    _, path, tokens = _render(RecipeId.WORK_REPORT, pages, tmp_path, "work")
    review = review_volume(path, pages, tmp_path, layout_scheme=tokens.layout_scheme)
    assert "card_soup" not in {issue.code for issue in review.issues}


def test_same_recipe_body_pages_are_not_clones(tmp_path: Path) -> None:
    pages = [
        _overview(1, ["增长换挡而非失速", "盈利底部渐明", "风格再平衡"], kpis=True),
        _expansion(2),
        _overview(3, ["信用扩张", "汇率稳定"]),
        _context(4),
    ]
    for recipe in (RecipeId.CONSULTING, RecipeId.WORK_REPORT, RecipeId.ART_MARKET, RecipeId.HISTORY):
        presentation, _, _ = _render(recipe, pages, tmp_path, f"mix-{recipe.value}")
        fingerprints = [_silhouette(slide) for slide in presentation.slides]
        assert len(set(fingerprints)) == len(fingerprints), (recipe.value, fingerprints)
