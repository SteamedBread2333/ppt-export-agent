from __future__ import annotations

import functools
import platform
import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE

from ppt_expert.models import (
    DesignSpec,
    OutlinePlan,
    StoryPage,
    ValidationIssue,
    ValidationReport,
)


def validate_presentation(
    pptx_path: str | Path,
    outline: OutlinePlan,
    story: list[StoryPage],
    design: DesignSpec,
    image_paths: dict[str, str],
) -> ValidationReport:
    path = Path(pptx_path).expanduser().resolve()
    issues: list[ValidationIssue] = []
    if not path.exists():
        return ValidationReport(
            valid=False,
            pptx_path=str(path),
            issues=[ValidationIssue(code="missing_pptx", message="PPTX file does not exist")],
        )

    presentation = Presentation(path)
    if len(presentation.slides) != len(outline.pages):
        issues.append(
            ValidationIssue(
                code="page_count",
                message=(
                    f"Expected {len(outline.pages)} slides; "
                    f"found {len(presentation.slides)}"
                ),
            )
        )
    if len(story) != len(outline.pages):
        issues.append(
            ValidationIssue(
                code="story_count",
                message="STORY slide count does not match the user outline",
            )
        )

    for index, outline_page in enumerate(outline.pages):
        if index >= len(story):
            break
        story_page = story[index]
        if story_page.number != outline_page.number or story_page.title != outline_page.title:
            issues.append(
                ValidationIssue(
                    code="outline_fidelity",
                    message=f"Slide number or title diverges from outline: {story_page.title}",
                    page=index + 1,
                )
            )
        if sum(len(item) for item in story_page.content) > 320:
            issues.append(
                ValidationIssue(
                    code="text_overflow_risk",
                    message="Body copy exceeds 320 characters and may overflow",
                    page=index + 1,
                )
            )
        if story_page.image_id:
            image_path = image_paths.get(story_page.image_id)
            if not image_path or not Path(image_path).exists():
                issues.append(
                    ValidationIssue(
                        code="missing_image",
                        message=f"Missing artwork: {story_page.image_id}",
                        page=index + 1,
                    )
                )

    for slide_index, slide in enumerate(presentation.slides, 1):
        page = story[slide_index - 1] if slide_index <= len(story) else None
        slide_text = "\n".join(
            shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        )
        if page and page.title not in slide_text:
            issues.append(
                ValidationIssue(
                    code="missing_title",
                    message=f"Rendered slide is missing its title: {page.title}",
                    page=slide_index,
                )
            )
        if page:
            for item in page.content:
                if item not in slide_text:
                    issues.append(
                        ValidationIssue(
                            code="missing_content",
                            message=f"Rendered slide is missing core content: {item}",
                            page=slide_index,
                        )
                    )
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > presentation.slide_width
                or shape.top + shape.height > presentation.slide_height
            ):
                issues.append(
                    ValidationIssue(
                        code="out_of_bounds",
                        message="An element extends beyond the slide boundary",
                        page=slide_index,
                    )
                )
                break

    palette = design.palette_hex()
    allowed_fonts = design.allowed_font_names()
    for slide_index, slide in enumerate(presentation.slides, 1):
        unexpected_colors = _shape_colors(slide) - palette
        if unexpected_colors:
            issues.append(
                ValidationIssue(
                    code="palette_violation",
                    message=(
                        "Found colors outside the approved palette: "
                        f"{', '.join(sorted(unexpected_colors))}"
                    ),
                    page=slide_index,
                )
            )
        used_fonts = _shape_fonts(slide)
        unexpected_fonts = {font for font in used_fonts if font.casefold() not in allowed_fonts}
        if unexpected_fonts:
            issues.append(
                ValidationIssue(
                    code="font_violation",
                    message=(
                        "Found fonts outside the DESIGN contract: "
                        f"{', '.join(sorted(unexpected_fonts))}"
                    ),
                    page=slide_index,
                )
            )

    font_chains = [
        [design.title_font, *design.title_font_fallbacks],
        [design.body_font, *design.body_font_fallbacks],
    ]
    for chain in font_chains:
        if not any(_font_installed(font) for font in chain):
            issues.append(
                ValidationIssue(
                    code="font_unavailable",
                    message=f"Font fallback chain is unavailable: {' → '.join(chain)}",
                    severity="warning",
                )
            )
    if len(palette) < 4:
        issues.append(
            ValidationIssue(
                code="weak_palette",
                message="The approved palette has insufficient visual separation",
                severity="warning",
            )
        )
    return ValidationReport(
        valid=not any(issue.severity == "error" for issue in issues),
        issues=issues,
        pptx_path=str(path),
    )


def _shape_colors(slide) -> set[str]:
    colors: set[str] = set()
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is not None and fill.type == MSO_FILL_TYPE.SOLID:
            color = fill.fore_color
            if color.type == MSO_COLOR_TYPE.RGB:
                colors.add(f"#{color.rgb}".upper())
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.font.color.type == MSO_COLOR_TYPE.RGB:
                colors.add(f"#{paragraph.font.color.rgb}".upper())
            for run in paragraph.runs:
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    colors.add(f"#{run.font.color.rgb}".upper())
    return colors


def _shape_fonts(slide) -> set[str]:
    fonts: set[str] = set()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.font.name:
                fonts.add(paragraph.font.name)
            fonts.update(run.font.name for run in paragraph.runs if run.font.name)
    return fonts


@functools.lru_cache(maxsize=32)
def _font_installed(font_family: str) -> bool:
    normalized = _font_key(font_family)
    fc_list = shutil.which("fc-list")
    if fc_list:
        result = subprocess.run(
            [fc_list, ":", "family"],
            check=False,
            capture_output=True,
            text=True,
            timeout=15,
        )
        if result.returncode == 0 and normalized in _font_key(result.stdout):
            return True
    if platform.system() == "Darwin" and normalized in _font_key(_macos_font_catalog()):
        return True
    font_dirs = [
        Path.home() / "Library/Fonts",
        Path("/Library/Fonts"),
        Path("/System/Library/Fonts"),
        Path("C:/Windows/Fonts"),
        Path("/usr/share/fonts"),
        Path("/usr/local/share/fonts"),
    ]
    return any(
        normalized in _font_key(path.stem)
        for directory in font_dirs
        if directory.exists()
        for path in directory.rglob("*")
        if path.is_file()
    )


@functools.lru_cache(maxsize=1)
def _macos_font_catalog() -> str:
    profiler = shutil.which("system_profiler")
    if profiler is None:
        return ""
    result = subprocess.run(
        [profiler, "SPFontsDataType", "-json"],
        check=False,
        capture_output=True,
        text=True,
        timeout=30,
    )
    return result.stdout if result.returncode == 0 else ""


def _font_key(value: str) -> str:
    return "".join(character for character in value.casefold() if character.isalnum()).removesuffix(
        "sc"
    )


def write_validation_report(report: ValidationReport, project_dir: Path) -> str:
    json_path = project_dir / "validation.json"
    md_path = project_dir / "VALIDATION.md"
    json_path.write_text(report.model_dump_json(indent=2), encoding="utf-8")
    lines = [
        "# VALIDATION",
        "",
        f"- Result: {'Passed' if report.valid else 'Failed'}",
        f"- File: {report.pptx_path}",
        "",
        "## Issues",
    ]
    if report.issues:
        for issue in report.issues:
            page = f"Slide {issue.page}: " if issue.page else ""
            lines.append(f"- [{issue.severity}] {page}{issue.message} (`{issue.code}`)")
    else:
        lines.append("- None")
    md_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return str(md_path.resolve())
