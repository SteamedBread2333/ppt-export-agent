from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation

from ppt_expert.models import GuardReport, GuardWarning

_SHORT = re.compile(
    r"(?:¥[\d,.]+[万亿MBK]?|\d+(?:\.\d+)?%|\d{4}\s*Q[1-4]|W\d{1,2}|20\d{2}[-/.]\d{1,2}(?:[-/.]\d{1,2})?)"
)


def inspect_guards(pptx_path: str | Path) -> GuardReport:
    presentation = Presentation(str(pptx_path))
    warnings: list[GuardWarning] = []
    for index, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text or ""
            width = shape.width.inches
            wrap = bool(shape.text_frame.word_wrap)
            for match in _SHORT.finditer(text):
                token = match.group(0)
                recommended = max(0.42, len(token) * 0.11)
                compact = text.strip() == token or len(text.strip()) <= len(token) + 2
                if wrap and compact and width + 1e-6 < recommended:
                    warnings.append(
                        GuardWarning(
                            page=index,
                            token=token,
                            box_width=round(width, 3),
                            recommended=round(recommended, 3),
                            message=(
                                f'"{token}" may still be too narrow: '
                                f"w={width:.2f}in, recommended={recommended:.2f}in"
                            ),
                        )
                    )
    return GuardReport(warnings=warnings)


def write_guard_report(report: GuardReport, project_dir: str | Path) -> str:
    path = Path(project_dir) / "guards.json"
    path.write_text(json.dumps(report.model_dump(mode="json"), indent=2, ensure_ascii=False), encoding="utf-8")
    return str(path.resolve())
