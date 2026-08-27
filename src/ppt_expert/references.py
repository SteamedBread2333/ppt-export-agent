from __future__ import annotations

from collections import Counter
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE

from ppt_expert.models import ReferenceAnalysis, StyleOption

DEFAULT_COLORS = ["#325D79", "#E8A87C", "#F7F4EF", "#1F2933", "#4FA3A5"]


def analyze_references(
    template_path: str | None,
    image_paths: list[str],
) -> ReferenceAnalysis:
    colors: Counter[str] = Counter()
    fonts: Counter[str] = Counter()
    notes: list[str] = []
    resolved_images = [str(Path(path).expanduser().resolve()) for path in image_paths]
    resolved_template = (
        str(Path(template_path).expanduser().resolve()) if template_path else None
    )

    layout_families: list[str] = []
    if resolved_template:
        template_colors, template_fonts, layout_families = _template_tokens(
            Path(resolved_template)
        )
        colors.update(template_colors)
        fonts.update(template_fonts)
        notes.append("Extracted colors, fonts, masters, and page dimensions from template")
        notes.append(
            "Reusable layout families: " + ", ".join(layout_families[:8] or ["blank"])
        )
    for image_path in resolved_images:
        try:
            colors.update(_image_palette(Path(image_path)))
        except OSError as exc:
            notes.append(f"Unable to read reference image {Path(image_path).name}: {exc}")

    palette = _semantic_palette(list(colors.elements()) or DEFAULT_COLORS)
    title_font = fonts.most_common(1)[0][0] if fonts else None
    body_font = fonts.most_common(2)[-1][0] if fonts else title_font
    source_type = "mixed" if resolved_template and resolved_images else (
        "template" if resolved_template else "images"
    )
    style = StyleOption(
        key="A",
        name="Reference-led direction",
        mood="Derived directly from user-linked templates and images",
        primary=palette[0],
        secondary=palette[1],
        background=palette[2],
        text=palette[3],
        accent=palette[4],
    )
    return ReferenceAnalysis(
        source_type=source_type,
        template_path=resolved_template,
        image_paths=resolved_images,
        preview_paths=[path for path in [resolved_template, *resolved_images] if path],
        style=style,
        title_font=title_font,
        body_font=body_font,
        notes=notes,
        layout_families=layout_families,
    )


def _template_tokens(path: Path) -> tuple[Counter[str], Counter[str], list[str]]:
    presentation = Presentation(path)
    colors: Counter[str] = Counter()
    fonts: Counter[str] = Counter()
    layout_families = [layout.name for layout in presentation.slide_layouts if layout.name]
    for slide in presentation.slides:
        for shape in slide.shapes:
            fill = getattr(shape, "fill", None)
            if fill is not None and fill.type == MSO_FILL_TYPE.SOLID:
                color = fill.fore_color
                if color.type == MSO_COLOR_TYPE.RGB:
                    colors[f"#{color.rgb}".upper()] += 3
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font.name:
                    fonts[paragraph.font.name] += 1
                if paragraph.font.color.type == MSO_COLOR_TYPE.RGB:
                    colors[f"#{paragraph.font.color.rgb}".upper()] += 1
                for run in paragraph.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1
                    if run.font.color.type == MSO_COLOR_TYPE.RGB:
                        colors[f"#{run.font.color.rgb}".upper()] += 1
    return colors, fonts, layout_families


def _image_palette(path: Path) -> Counter[str]:
    with Image.open(path) as image:
        sample = image.convert("RGB")
        sample.thumbnail((256, 256))
        quantized = sample.quantize(colors=8, method=Image.Quantize.MEDIANCUT)
        raw_palette = quantized.getpalette() or []
        counts = Counter(quantized.get_flattened_data())
    result: Counter[str] = Counter()
    for index, count in counts.items():
        offset = index * 3
        red, green, blue = raw_palette[offset : offset + 3]
        result[f"#{red:02X}{green:02X}{blue:02X}"] += count
    return result


def _semantic_palette(colors: list[str]) -> list[str]:
    unique = list(dict.fromkeys(color.upper() for color in colors))
    for fallback in DEFAULT_COLORS:
        if fallback not in unique:
            unique.append(fallback)
    background = max(unique, key=_luminance)
    text = min(unique, key=_luminance)
    candidates = [color for color in unique if color not in {background, text}]
    candidates.sort(key=_saturation, reverse=True)
    primary, secondary, accent = (candidates + DEFAULT_COLORS)[:3]
    return [primary, secondary, background, text, accent]


def _channels(color: str) -> tuple[int, int, int]:
    value = color.lstrip("#")
    return int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def _luminance(color: str) -> float:
    red, green, blue = _channels(color)
    return 0.2126 * red + 0.7152 * green + 0.0722 * blue


def _saturation(color: str) -> int:
    channels = _channels(color)
    return max(channels) - min(channels)
