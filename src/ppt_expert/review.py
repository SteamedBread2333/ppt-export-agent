from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ppt_expert.models import LayoutScheme, PageRole, QualityIssue, StoryPage, VolumeReview
from ppt_expert.preview import render_montage


def review_volume(
    pptx_path: str | Path,
    pages: list[StoryPage],
    project_dir: str | Path,
    *,
    dpi: int = 70,
    layout_scheme: LayoutScheme = LayoutScheme.RULES,
    native_edit: bool = False,
) -> VolumeReview:
    presentation = Presentation(str(pptx_path))
    issues: list[QualityIssue] = []
    roles = [page.resolved_role() for page in pages]
    rhythm_ok = bool(pages) and pages[0].is_dark() and pages[-1].is_dark()
    if pages and not rhythm_ok:
        issues.append(
            QualityIssue(
                code="weak_dark_bookends",
                message="Cover and close should use the dark master for volume rhythm",
                cause="LIGHT master was used on an identity page",
                repair_scope="rhythm",
                acceptance="Cover and close render on DARK",
            )
        )
    adjacent_repeat = False
    for index in range(1, len(roles)):
        if roles[index] == roles[index - 1] and roles[index] not in {PageRole.EVIDENCE}:
            adjacent_repeat = True
            issues.append(
                QualityIssue(
                    code="adjacent_silhouette",
                    message="Adjacent slides share the same narrative silhouette",
                    page=index + 1,
                    cause="Page roles were not varied",
                    repair_scope="rhythm",
                    acceptance="Neighboring slides use different roles",
                )
            )
            break
    empty_bottom = False
    chrome_locked = True
    footer_ys: list[float] = []
    if not native_edit:
        for slide_index, slide in enumerate(presentation.slides):
            bottoms = []
            for shape in slide.shapes:
                top = shape.top.inches
                height = shape.height.inches
                bottoms.append(top + height)
                if 7.0 <= top <= 7.3:
                    footer_ys.append(round(top, 2))
            if bottoms and max(bottoms) < 5.4:
                empty_bottom = True
                issues.append(
                    QualityIssue(
                        code="empty_bottom",
                        message="The lower third of the slide is unintentionally empty",
                        page=slide_index + 1,
                        cause="Evidence band does not reach 88–92%",
                        repair_scope="slide",
                        acceptance="Content or implication occupies the lower band",
                    )
                )
        if footer_ys and len(set(footer_ys)) > 2:
            chrome_locked = False
            issues.append(
                QualityIssue(
                    code="footer_drift",
                    message="Footer band is not locked across the deck",
                    cause="Folio y drifted between slides",
                    repair_scope="token",
                    acceptance="Footer y is shared",
                )
            )
    representative = _representatives(pages)
    montage, pdf = render_montage(pptx_path, Path(project_dir) / "render", dpi=dpi)
    issues.extend(
        inspect_representatives(pptx_path, pages, representative, native_edit=native_edit)
    )
    issues.extend(_numeric_conflicts(pages))
    if not native_edit:
        issues.extend(_aesthetic_geometry(presentation, pages, layout_scheme))
    return VolumeReview(
        rhythm_ok=rhythm_ok,
        no_adjacent_repeat=not adjacent_repeat,
        no_empty_bottom=not empty_bottom,
        chrome_locked=chrome_locked,
        representative_pages=representative,
        issues=issues,
        montage_path=montage,
        pdf_path=pdf,
    )


def _representatives(pages: list[StoryPage]) -> list[int]:
    if not pages:
        return []
    densest = max(pages, key=lambda page: sum(len(item) for item in page.content))
    special = next(
        (page for page in pages if page.chart or page.table or page.allocation or page.scenarios),
        densest,
    )
    chosen = [pages[0].number, densest.number, special.number, pages[-1].number]
    return list(dict.fromkeys(chosen))


def inspect_representatives(
    pptx_path: str | Path,
    pages: list[StoryPage],
    numbers: list[int],
    *,
    native_edit: bool = False,
) -> list[QualityIssue]:
    """Four questions on representative pages: assertion, overflow, columns, implication."""
    issues: list[QualityIssue] = []
    by_number = {page.number: page for page in pages}
    presentation = Presentation(str(pptx_path))
    for number in numbers:
        page = by_number.get(number)
        if page is None or number > len(presentation.slides):
            continue
        if not _is_assertion_title(page.title):
            issues.append(
                QualityIssue(
                    code="topic_title",
                    message="Representative title reads as a topic, not an assertion",
                    page=number,
                    severity="error",
                    cause="The title names a subject instead of a judgment",
                    repair_scope="slide",
                    acceptance="Title is a complete assertion",
                )
            )
        slide = presentation.slides[number - 1]
        xs: list[float] = []
        implication_top = None
        for shape in slide.shapes:
            left = shape.left
            width = shape.width
            top = shape.top
            height = shape.height
            if (
                left < 0
                or top < 0
                or left + width > presentation.slide_width
                or top + height > presentation.slide_height
            ):
                issues.append(
                    QualityIssue(
                        code="rep_overflow",
                        message="A shape on a representative page overflows the canvas",
                        page=number,
                        cause="Element extends past the slide edge",
                        repair_scope="slide",
                        acceptance="All shapes sit inside the slide canvas",
                    )
                )
                break
            if native_edit:
                continue
            if (
                getattr(shape, "has_text_frame", False)
                and page.takeaway
                and page.takeaway[:10] in (shape.text or "")
                and 6.1 <= shape.top.inches <= 6.75
            ):
                implication_top = shape.top.inches
            if 1.1 <= shape.top.inches <= 1.3 and shape.width.inches > 1.5:
                xs.append(round(shape.left.inches, 2))
        if native_edit:
            continue
        if implication_top is not None and not (6.2 <= implication_top <= 6.55):
            issues.append(
                QualityIssue(
                    code="implication_drift",
                    message="Conclusion bar is not hugging the folio safety band",
                    page=number,
                    cause="Implication y drifted off the locked token",
                    repair_scope="token",
                    acceptance="Implication sits at the locked y near the footer",
                )
            )
        if len(xs) >= 2:
            span = max(xs) - min(xs)
            step = span / max(len(set(xs)) - 1, 1)
            if step and any(abs((x - min(xs)) % step) > 0.12 and abs(((x - min(xs)) % step) - step) > 0.12 for x in xs):
                issues.append(
                    QualityIssue(
                        code="column_misalign",
                        message="Column modules on a representative page are not on a shared grid",
                        page=number,
                        cause="Left edges do not share a repeating interval",
                        repair_scope="slide",
                        acceptance="Sibling columns share a left-edge cadence",
                    )
                )
    return issues


def _is_assertion_title(title: str) -> bool:
    stripped = title.strip()
    if any(mark in stripped for mark in "，,;；:：.—–"):
        return True
    return len(stripped) >= 12


def _numeric_conflicts(pages: list[StoryPage]) -> list[QualityIssue]:
    seen: dict[str, str] = {}
    issues: list[QualityIssue] = []
    for page in pages:
        for item in page.kpis:
            key = item.label.strip().casefold()
            if not key:
                continue
            if key in seen and seen[key] != item.value:
                issues.append(
                    QualityIssue(
                        code="conflicting_kpi",
                        message=f'KPI "{item.label}" is quoted with conflicting values',
                        page=page.number,
                        cause="The same metric is restated with a different number",
                        repair_scope="narrative",
                        acceptance="A named metric keeps one value across the deck",
                    )
                )
                break
            seen[key] = item.value
    return issues


def _aesthetic_geometry(
    presentation, pages: list[StoryPage], layout_scheme: LayoutScheme
) -> list[QualityIssue]:
    """Taste the montage can enforce without a vision model."""
    issues: list[QualityIssue] = []
    slide_w = presentation.slide_width
    slide_h = presentation.slide_height
    for index, page in enumerate(pages):
        if index >= len(presentation.slides):
            break
        slide = presentation.slides[index]
        role = page.resolved_role()
        if (
            layout_scheme in {LayoutScheme.RULES, LayoutScheme.SPREAD}
            and role not in {PageRole.COVER, PageRole.CLOSE, PageRole.SCENARIO}
        ):
            cards = _filled_modules(slide, slide_w, slide_h)
            if cards >= 2:
                issues.append(
                    QualityIssue(
                        code="card_soup",
                        message="The slide is a stack of filled tiles instead of type, rules, and numbers",
                        page=page.number,
                        severity="error",
                        cause="Compose wrapped evidence in surface cards",
                        repair_scope="slide",
                        acceptance="Body pages use hairlines and hierarchy, not gray modules",
                    )
                )
        if role not in {PageRole.COVER, PageRole.CLOSE} and _title_is_cramped(slide, page):
            issues.append(
                QualityIssue(
                    code="cramped_header",
                    message="Title band collides with the evidence; there is no pause after the rule",
                    page=page.number,
                    severity="error",
                    cause="content_top sits on the header hairline",
                    repair_scope="token",
                    acceptance="At least 0.22in of air between the title box and the first content",
                )
            )
    return issues


def _filled_modules(slide, slide_w, slide_h) -> int:
    from pptx.enum.dml import MSO_FILL_TYPE

    count = 0
    page_w = slide_w.inches
    page_h = slide_h.inches
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
            continue
        width = shape.width.inches
        height = shape.height.inches
        if width >= page_w - 0.08 and height >= page_h - 0.08:
            continue
        if height < 0.5 or width < 1.6:
            continue
        fill = getattr(shape, "fill", None)
        try:
            if fill is None or fill.type != MSO_FILL_TYPE.SOLID:
                continue
        except (AttributeError, TypeError, ValueError):
            continue
        count += 1
    return count


def _title_is_cramped(slide, page: StoryPage) -> bool:
    title_bottom = None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and page.title in (shape.text or ""):
            title_bottom = shape.top.inches + min(shape.height.inches, 0.55)
            break
    if title_bottom is None:
        return False
    content_tops = []
    for shape in slide.shapes:
        top = shape.top.inches
        height = shape.height.inches
        if height < 0.05 or top <= title_bottom + 0.02 or top >= 6.9:
            continue
        content_tops.append(top)
    if not content_tops:
        return False
    return min(content_tops) - title_bottom < 0.22
